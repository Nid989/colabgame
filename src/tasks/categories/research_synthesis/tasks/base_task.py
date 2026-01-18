"""
Base generator class for Information Synthesis & Presentation tasks.
"""

import random
import string
import os
import json
from typing import Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches

from ...base import BaseTask
from ..enhanced_content_generator import SimplifiedContentGenerator


class BaseResearchSynthesisTaskGenerator(BaseTask):
    """Base class for Information Synthesis & Presentation task implementations."""

    def __init__(self, task_type: str, level: int):
        super().__init__()
        self.task_type = task_type
        self.level = level
        self.random = random.Random()

        # Initialize enhanced content generator
        self.content_generator = SimplifiedContentGenerator()

        # Available domains for content generation
        self.domains = list(self.content_generator.domains.keys())

        # Content types by level
        self.content_types = {1: ["company_info", "product_info", "event_info"], 2: ["product_catalog", "contact_info"], 3: ["document_repository"]}

    def set_seed(self, seed: int):
        """Set random seed for reproducible generation."""
        self.random.seed(seed)

    def generate(self, **kwargs) -> Dict[str, Any]:
        """Generate a task instance (implemented by calling generate_task_data)."""
        seed = kwargs.get("seed")
        return self.generate_task_data(seed)

    def generate_task_data(self, seed: Optional[int] = None) -> Dict[str, Any]:
        """Generate dynamic task data."""
        return self.generate_basic_task_structure(seed)

    def get_task_type(self) -> str:
        """Return the task type identifier."""
        return self.task_type

    def get_level(self) -> int:
        """Return the difficulty level of this task."""
        return self.level

    def generate_random_string(self, length: int = 8) -> str:
        """Generate random string for identifiers."""
        return "".join(self.random.choices(string.ascii_lowercase, k=length))

    def generate_random_number(self, min_val: int = 1000, max_val: int = 9999) -> int:
        """Generate random number for IDs."""
        return self.random.randint(min_val, max_val)

    def generate_task_id(self) -> str:
        """Generate unique task identifier."""
        return f"{self.task_type}_{self.generate_random_number()}"

    def generate_presentation_filename(self) -> str:
        """Generate presentation filename."""
        task_id = self.generate_random_number()
        return f"presentation_{task_id}.pptx"

    def generate_webpage_content(self, level: int, seed: Optional[int] = None) -> Dict[str, Any]:
        """Generate enhanced webpage content based on level and seed using domain-based approach."""
        if seed is not None:
            self.set_seed(seed)

        # Select domain and content type
        domain = self.random.choice(self.domains)
        available_content_types = self.content_types.get(level, [])

        if not available_content_types:
            raise ValueError(f"No content types available for level {level}")

        content_type = self.random.choice(available_content_types)

        # Use the enhanced content generator
        if level == 1:
            return self.content_generator.generate_level1_content(domain, content_type, seed)
        elif level == 2:
            return self.content_generator.generate_level2_content(domain, content_type, seed)
        elif level == 3:
            return self.content_generator.generate_level3_content(domain, seed)
        else:
            raise ValueError(f"Unsupported level: {level}")

    def create_expected_presentation(self, task_data: Dict[str, Any], temp_dir: str) -> str:
        """Create expected presentation file for evaluation."""
        prs = Presentation()

        if self.level == 1:
            # Single slide with target text
            slide = prs.slides.add_slide(prs.slide_layouts[5])  # Blank layout
            textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
            textbox.text_frame.text = task_data["target_text"]

        elif self.level == 2:
            # Two slides with different texts
            for i, text in enumerate(task_data["target_texts"]):
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
                textbox.text_frame.text = text

        elif self.level == 3:
            # Single slide with file content
            slide = prs.slides.add_slide(prs.slide_layouts[5])
            textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(4))
            textbox.text_frame.text = task_data["file_content"]

        filename = f"expected_{task_data['presentation_file']}"
        filepath = os.path.join(temp_dir, filename)
        prs.save(filepath)
        return filepath

    def create_expected_spec(self, task_data: Dict[str, Any], temp_dir: str) -> str:
        """Create expected presentation specification JSON file for evaluation."""

        if self.level == 1:
            # Level 1: Global required texts (user feedback)
            spec = {"global": {"required_texts": [task_data["target_text"]]}, "options": {"text_threshold": 0.8}}

        elif self.level == 2:
            # Level 2: Two slides with by_index matching
            spec = {
                "slides": [
                    {"match": {"by_index": 0}, "required": {"texts_all": [task_data["target_texts"][0]]}},
                    {"match": {"by_index": 1}, "required": {"texts_all": [task_data["target_texts"][1]]}},
                ],
                "options": {"text_threshold": 0.8},
            }

        elif self.level == 3:
            # Level 3: Single slide with file content
            spec = {"slides": [{"required": {"texts_all": [task_data["file_content"]]}}], "options": {"text_threshold": 0.8}}
        else:
            raise ValueError(f"Unsupported level: {self.level}")

        # Create JSON filename with same naming pattern but .json extension
        presentation_base = task_data["presentation_file"].replace(".pptx", "")
        filename = f"expected_{presentation_base}.json"
        filepath = os.path.join(temp_dir, filename)

        with open(filepath, "w", encoding="utf-8") as f:
            json.dump(spec, f, indent=2, ensure_ascii=False)

        return filepath

    def generate_basic_task_structure(self, seed: Optional[int] = None) -> Dict[str, Any]:
        """Generate basic task structure common to all tasks."""
        if seed is not None:
            self.set_seed(seed)

        return {
            "task_type": self.task_type,
            "level": self.level,
            "category": "research_synthesis",
            "seed": seed or self.random.randint(1, 10000),
            "presentation_file": self.generate_presentation_filename(),
            "evaluation_mode": "multi_evaluator",
            "example_id": f"L{self.level}_{self.task_type}_{self.random.randint(1, 1000)}",
            "file_name": self.generate_presentation_filename(),  # Required by framework
        }
